import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide():
    prs = Presentation()
    # Use a blank layout (usually index 6 in standard templates)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Helper function to add shapes with text
    def add_shape(shape_type, left, top, width, height, text, bg_color=None, font_size=10, bold=False, shape_name=""):
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.text = text
        shape.name = shape_name
        
        # Style
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        
        # Border (Black, thin)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(1)
        
        # Text Style
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(font_size)
        paragraph.font.name = 'Arial'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.font.bold = bold
        
        return shape

    # Helper function for connectors
    def add_connector(shape_from, shape_to, line_type="Straight"):
        # Simple straight line connector logic (can be adjusted for bent lines)
        # For this script, we'll manually draw lines for better control or use simple connectors
        connector = slide.shapes.add_connector(
            1, # msoConnectorStraight
            shape_from.left + shape_from.width/2,
            shape_from.top + shape_from.height,
            shape_to.left + shape_to.width/2,
            shape_to.top
        )
        connector.line.color.rgb = RGBColor(0, 0, 0)
        connector.line.width = Pt(1)
        # Add arrow head
        connector.line.end_arrowhead_style = 2 # Triangle
        return connector

    # Helper for manual arrows (more precise for complex diagrams)
    def add_arrow(x1, y1, x2, y2):
        line = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.end_arrowhead_style = 2
        return line

    # --- Colors (approximated from image) ---
    c_blue = RGBColor(189, 215, 238)
    c_yellow = RGBColor(255, 242, 204) # Forced Prompt
    c_green = RGBColor(198, 224, 180) # Probs / Student PRM
    c_purple = RGBColor(204, 192, 218) # Ensemble
    c_red = RGBColor(244, 177, 131)    # Reward / Scores
    c_grey = RGBColor(217, 217, 217)   # Feature Extraction
    c_orange = RGBColor(252, 228, 214) # GBDT
    c_final_blue = RGBColor(189, 215, 238) # Final Circle

    # --- Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Figure 1: Overview of the Proposed Rollout-free PRM Framework"
    p.font.bold = True
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # --- Vertical Separator ---
    sep = slide.shapes.add_connector(1, Inches(5), Inches(1), Inches(5), Inches(7))
    sep.line.color.rgb = RGBColor(0,0,0)
    sep.line.width = Pt(1.5)

    # ==========================================
    # (a) Training Section (Left)
    # ==========================================
    
    # Header
    h1 = slide.shapes.add_textbox(Inches(0.2), Inches(1), Inches(4.5), Inches(0.4))
    h1.text_frame.text = "(a) Training: Prompted Forcing & Ensemble Distillation"
    h1.text_frame.paragraphs[0].font.bold = True
    
    # Top Flow: Problem -> Steps
    prob_tr = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.2, 1.5, 1.0, 0.4, "Problem x", c_blue)
    step1 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 1.6, 1.5, 0.7, 0.4, "Step s1", RGBColor(240,240,240))
    step2 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 2.6, 1.5, 0.7, 0.4, "Step s2", RGBColor(240,240,240))
    dots = slide.shapes.add_textbox(Inches(3.4), Inches(1.5), Inches(0.5), Inches(0.4))
    dots.text_frame.text = "..."
    
    step_t = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 2.6, 2.3, 0.7, 0.4, "Step st", RGBColor(240,240,240))

    # Forced Prompt
    prompt = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.2, 2.3, 2.0, 0.6, "Forced Prompt\n(\"The final answer is []\")", c_yellow)
    # Add glow manually is hard in basic pptx, skipping glow but color matches.
    
    # Connection logic for training top
    add_arrow(1.2, 1.7, 1.6, 1.7) # Prob -> s1
    add_arrow(2.3, 1.7, 2.6, 1.7) # s1 -> s2
    add_arrow(3.3, 1.7, 3.5, 1.7) # s2 -> ...
    
    # Arrows merging to step t
    # s1 down
    line_s1 = slide.shapes.add_connector(1, Inches(1.95), Inches(1.9), Inches(1.95), Inches(2.1))
    line_s1.line.color.rgb = RGBColor(0,0,0)
    # s2 down
    line_s2 = slide.shapes.add_connector(1, Inches(2.95), Inches(1.9), Inches(2.95), Inches(2.3)) # Connects to st
    line_s2.line.color.rgb = RGBColor(0,0,0)
    line_s2.line.end_arrowhead_style = 2

    # Prompt to Step t integration (circle representation)
    circle_plus = add_shape(MSO_SHAPE.OVAL, 2.3, 2.4, 0.2, 0.2, "+", RGBColor(255,255,255))
    add_arrow(2.2, 2.5, 2.3, 2.5) # Prompt -> Circle
    
    # Teacher Models
    tm1 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.2, 3.3, 1.4, 0.8, "Teacher\nModel M1\n(e.g., DeepSeek)", c_blue)
    tm2 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, 3.3, 1.4, 0.8, "Teacher\nModel M2\n(e.g., Llama)", c_blue)
    tm3 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 3.4, 3.3, 1.4, 0.8, "Teacher\nModel M3\n(e.g., Qwen)", c_blue)

    # Arrows to teachers
    # From Step t / Prompt mix
    add_arrow(2.5, 2.6, 2.5, 3.0) # Down from +/st
    # Split line
    line_split = slide.shapes.add_connector(1, Inches(0.9), Inches(3.0), Inches(4.1), Inches(3.0))
    line_split.line.color.rgb = RGBColor(0,0,0)
    add_arrow(0.9, 3.0, 0.9, 3.3) # To M1
    add_arrow(2.5, 3.0, 2.5, 3.3) # To M2
    add_arrow(4.1, 3.0, 4.1, 3.3) # To M3

    # "No Rollout Required" text
    nr_text = slide.shapes.add_textbox(Inches(3.8), Inches(2.7), Inches(1.5), Inches(0.5))
    nr_text.text_frame.text = "No Rollout\nRequired"
    nr_text.text_frame.paragraphs[0].font.bold = True
    nr_text.text_frame.paragraphs[0].font.size = Pt(9)
    # Arrow for text
    add_arrow(4.0, 2.9, 3.8, 3.1) # Pointing roughly to input area

    # Probabilities
    p1 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 4.5, 0.8, 0.3, "P1(yGT)", c_green)
    p2 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 2.1, 4.5, 0.8, 0.3, "P2(yGT)", c_green)
    p3 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 3.7, 4.5, 0.8, 0.3, "P3(yGT)", c_green)
    
    add_connector(tm1, p1)
    add_connector(tm2, p2)
    add_connector(tm3, p3)

    # Ensemble
    ens = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 5.1, 2.0, 0.6, "Ensemble Average\n& Log-Sigmoid", c_purple)
    
    # Arrows to Ensemble
    add_arrow(0.9, 4.8, 1.8, 5.1)
    add_arrow(2.5, 4.8, 2.5, 5.1)
    add_arrow(4.1, 4.8, 3.2, 5.1)

    # Dense Reward Label
    rew = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, 6.0, 1.4, 0.6, "Dense Reward\nLabel rt", c_red)
    add_connector(ens, rew)

    # Student PRM Training
    st_train = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 3.6, 6.0, 1.2, 0.8, "Student PRM\nTraining\n(MSE Loss)", RGBColor(240,240,240))
    
    # Dashed arrow
    dash_arrow = slide.shapes.add_connector(1, Inches(3.2), Inches(6.3), Inches(3.6), Inches(6.3))
    dash_arrow.line.color.rgb = RGBColor(0,0,0)
    dash_arrow.line.dash_style = 6 # Dash
    dash_arrow.line.end_arrowhead_style = 2

    # ==========================================
    # (b) Inference Section (Right)
    # ==========================================
    
    # Header
    h2 = slide.shapes.add_textbox(Inches(5.2), Inches(1), Inches(4.5), Inches(0.4))
    h2.text_frame.text = "(b) Inference: Learned Aggregation Reranking"
    h2.text_frame.paragraphs[0].font.bold = True

    # Problem
    prob_inf = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 5.3, 1.8, 1.0, 0.4, "Problem x", c_blue)
    
    # Text: Candidate Generation
    cg_text = slide.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(1.5), Inches(0.6))
    cg_text.text_frame.text = "Candidate\nGeneration\n(Best-of-N)"
    cg_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cg_text.text_frame.paragraphs[0].font.size = Pt(9)
    
    add_arrow(6.3, 2.0, 6.7, 2.0)

    # Fan out arrows
    add_arrow(7.8, 1.8, 8.2, 1.5)
    add_arrow(7.8, 1.9, 8.2, 1.7)
    add_arrow(7.8, 2.1, 8.2, 2.3)
    add_arrow(7.8, 2.2, 8.2, 2.5)

    # Student PRM (Inference)
    sprm_inf = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8.3, 1.6, 1.4, 0.7, "Student PRM\n(Inference)", c_green)

    # Scores v1, v2, vT
    score_v1 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 6.3, 3.2, 0.8, 0.4, "Score v1", c_red)
    score_v2 = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 7.3, 3.2, 0.8, 0.4, "Score v2", c_red)
    
    score_dots = slide.shapes.add_textbox(Inches(8.15), Inches(3.2), Inches(0.4), Inches(0.4))
    score_dots.text_frame.text = "..."
    
    score_vt = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8.6, 3.2, 0.8, 0.4, "Score vT", c_red)

    # Arrows from PRM to Scores (curved in image, straight here for simplicity or multi-segment)
    # To v1
    add_arrow(8.5, 2.3, 6.7, 3.2)
    # To v2
    add_arrow(8.7, 2.3, 7.7, 3.2)
    # To vT
    add_arrow(9.0, 2.3, 9.0, 3.2)

    # Arrows between scores
    add_arrow(7.1, 3.4, 7.3, 3.4)
    add_arrow(8.1, 3.4, 8.3, 3.4) # imaginary over dots
    add_arrow(8.3, 3.4, 8.6, 3.4)

    # Feature Extraction
    feat = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 5.2, 5.0, 1.2, 1.5, "Feature\nExtraction\n\n[Min, Mean,\nMax, Last-3\nMin, Sum-Logits,\n...]", c_grey)
    feat.text_frame.paragraphs[0].font.bold = True
    
    # Connect Scores to Feature Extraction
    # One big connector gathering them
    conn_gather = slide.shapes.add_connector(1, Inches(6.7), Inches(3.6), Inches(6.7), Inches(4.2)) # v1 down
    conn_gather.line.color.rgb = RGBColor(0,0,0)
    
    conn_gather2 = slide.shapes.add_connector(1, Inches(9.0), Inches(3.6), Inches(7.8), Inches(4.2)) # vT down-left
    conn_gather2.line.color.rgb = RGBColor(0,0,0)
    
    # Line to feature box
    add_arrow(6.7, 4.2, 5.8, 5.0) # From gathering point to feature box
    add_arrow(7.8, 4.2, 5.8, 5.0) # From gathering point 2

    # Learned Aggregation (Diamond)
    agg = add_shape(MSO_SHAPE.DIAMOND, 6.8, 5.1, 1.0, 1.3, "Learned\nAggregation\n(GBDT)", c_orange)
    
    add_arrow(6.4, 5.75, 6.8, 5.75) # Feat -> Agg

    # Final Path Score (Circle)
    fps = add_shape(MSO_SHAPE.OVAL, 8.0, 5.4, 0.7, 0.7, "Final\nPath\nScore", c_blue)
    add_arrow(7.8, 5.75, 8.0, 5.75)

    # Reranking & Selection
    rerank = add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 9.0, 5.5, 1.0, 0.6, "Reranking\n& Selection", RGBColor(240,240,240))
    add_arrow(8.7, 5.75, 9.0, 5.75)


    # Save
    prs.save('rollout_free_prm.pptx')
    print("Presentation saved as rollout_free_prm.pptx")

create_slide()